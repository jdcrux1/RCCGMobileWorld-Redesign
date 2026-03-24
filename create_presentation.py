#!/usr/bin/env python3
"""
RCCG Mobile World - ICT Board Presentation Generator
This script creates a PowerPoint presentation from the markdown content.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import sys

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
NAVY = RGBColor(15, 23, 42)      # #0F172A
BLUE = RGBColor(59, 130, 246)     # #3B82F6
GOLD = RGBColor(217, 119, 6)        # #D97706
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(156, 163, 175)
LIGHT_GRAY = RGBColor(248, 250, 252)

def add_title_slide(prs, title, subtitle, presenter, date):
    """Create title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background shape
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER
    
    # Presenter info
    pres_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12.333), Inches(1.5))
    tf = pres_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"Presented to: ICT Board of Directors\nDate: {date}\nBy: {presenter}"
    p.font.size = Pt(18)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, bullets, subtitle=None):
    """Create content slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = NAVY
    title_bar.line.fill.background()
    
    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(22)
        p.font.color.rgb = WHITE
        p.space_after = Pt(12)
    
    # Subtitle if provided
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12.333), Inches(0.5))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(14)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER

def add_two_column_slide(prs, title, left_title, left_bullets, right_title, right_bullets):
    """Create two column slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = NAVY
    title_bar.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Left column title
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(6), Inches(0.5))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = GOLD
    
    # Left content
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(6), Inches(4.5))
    tf = left_box.text_frame
    for i, bullet in enumerate(left_bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.space_after = Pt(8)
    
    # Right column title
    right_title_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.4), Inches(6), Inches(0.5))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = BLUE
    
    # Right content
    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(2), Inches(6), Inches(4.5))
    tf = right_box.text_frame
    for i, bullet in enumerate(right_bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.space_after = Pt(8)

def add_image_placeholder_slide(prs, title, description):
    """Create slide with image placeholder"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = NAVY
    title_bar.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Image placeholder box
    placeholder = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), Inches(2), Inches(9.333), Inches(4))
    placeholder.fill.solid()
    placeholder.fill.fore_color.rgb = RGBColor(30, 41, 59)
    placeholder.line.color.rgb = GRAY
    
    # Description
    desc_box = slide.shapes.add_textbox(Inches(2), Inches(4), Inches(9.333), Inches(1))
    tf = desc_box.text_frame
    p = tf.paragraphs[0]
    p.text = description
    p.font.size = Pt(18)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER

def add_timeline_slide(prs, title, quarters):
    """Create timeline slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = NAVY
    title_bar.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Timeline boxes
    colors = [NAVY, BLUE, GOLD, NAVY]
    for i, (quarter, items) in enumerate(quarters.items()):
        x = Inches(0.5 + i * 3.2)
        y = Inches(1.5)
        w = Inches(3)
        h = Inches(5.5)
        
        # Box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        box.fill.solid()
        box.fill.fore_color.rgb = colors[i]
        box.line.color.rgb = GRAY
        
        # Quarter title
        q_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.2), w - Inches(0.4), Inches(0.6))
        tf = q_box.text_frame
        p = tf.paragraphs[0]
        p.text = quarter
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        
        # Items
        items_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(1), w - Inches(0.4), Inches(4))
        tf = items_box.text_frame
        for j, item in enumerate(items):
            if j == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = "• " + item
            p.font.size = Pt(14)
            p.font.color.rgb = WHITE

def add_closing_slide(prs, title, message, contact_info):
    """Create closing slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Message
    msg_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12.333), Inches(2))
    tf = msg_box.text_frame
    p = tf.paragraphs[0]
    p.text = message
    p.font.size = Pt(24)
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER
    
    # Contact info
    contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12.333), Inches(1.5))
    tf = contact_box.text_frame
    p = tf.paragraphs[0]
    p.text = contact_info
    p.font.size = Pt(18)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER

# ============================================================================
# PRESENTATION CONTENT
# ============================================================================

print("Creating RCCG Mobile World ICT Board Presentation...")

# Slide 1: Title
add_title_slide(prs, 
    "RCCG MOBILE WORLD", 
    "Digital Transformation Initiative\nBridging Faith & Technology",
    "Your Name, Head of System Analysis & Frontend Design",
    "March 2026")

# Slide 2: Executive Summary
add_content_slide(prs, "Executive Summary", [
    "Vision: Create a unified digital platform for RCCG globally",
    "Current Status: Frontend Phase Complete",
    "Investment: Proposed budget for backend development",
    "ROI: Projected increase in engagement, giving, member retention",
    "Request: Approval for backend development phase"
])

# Slide 3: The Opportunity
add_two_column_slide(prs, "The Opportunity", 
    "Challenge:", [
        "50,000+ parishes worldwide",
        "Millions of members offline",
        "Inconsistent digital experiences",
        "Limited online giving infrastructure"
    ],
    "Solution:", [
        "Unified digital platform",
        "Mobile-first approach",
        "Standardized UX",
        "Secure payment processing"
    ])

# Slide 4: Timeline
add_timeline_slide(prs, "Project Timeline 2026", {
    "Q1 2026": ["Frontend Complete", "Backend Planning", "API Design", "Team Setup"],
    "Q2 2026": ["Backend Dev", "Auth System", "User Features", "Beta Launch"],
    "Q3 2026": ["Payment Integration", "Mobile App", "Analytics", "Full Launch"],
    "Q4 2026": ["AI Features", "Multi-language", "Scale & Grow", "2027 Planning"]
})

# Slide 5: Success Metrics
add_content_slide(prs, "Success Metrics", [
    "50% increase in page views - target",
    "30% growth in online giving",
    "60% reduction in admin time",
    "40% increase in mobile users",
    "User Satisfaction Target: 4.5/5",
    "Accessibility: WCAG 2.1 AA Compliant"
])

# Slide 6: Problem Statement
add_content_slide(prs, "Problem Statement", [
    "Fragmented digital presence across parishes",
    "No centralized content management",
    "Limited online giving options",
    "Poor mobile experience",
    "No member authentication system",
    "Manual event registration processes"
])

# Slide 7: User Personas
add_content_slide(prs, "User Personas", [
    "PASTOR: Preach, Manage, Connect - Needs easy content upload",
    "MEMBER: Worship, Learn, Give - Wants personalized experience",
    "VISITOR: Explore, Connect, Join - Seeks seamless navigation",
    "78% access via mobile devices",
    "65% want online giving options",
    "82% want sermon access anytime"
])

# Slide 8: Information Architecture
add_content_slide(prs, "Information Architecture", [
    "HOME - Hero, Latest Sermon, Upcoming Events, Give CTA",
    "ABOUT - Story, Beliefs, Staff, Ministries",
    "SERMONS - Latest, Series, Topics, Audio/Video",
    "EVENTS - Upcoming, Past, Registration",
    "MINISTRIES - Youth, Women, Children, Groups",
    "GIVE - Online, Recurring, Goals, History (Auth required)",
    "Authenticated: Dashboard, Profile, Bookmarks"
])

# Slide 9: Wireframes (placeholder)
add_image_placeholder_slide(prs, "Home Page Wireframe", 
    "[Insert wireframe image here]\nMobile-first responsive design with hero section, featured content, and navigation")

# Slide 10: Technical Architecture
add_content_slide(prs, "Technical Architecture", [
    "Frontend: HTML5, CSS3, JavaScript, Tailwind CSS, Vercel",
    "Backend (to be): Node.js/Express, PostgreSQL, Redis",
    "API Gateway: Authentication, Rate Limiting",
    "Services: User, Content, Payment microservices",
    "Storage: AWS S3 for media, CloudFlare CDN",
    "Security: JWT, HTTPS, OWASP compliance"
])

# Slide 11: Database Schema
add_content_slide(prs, "Database Schema Overview", [
    "USERS: id, email, password_hash, name, role, created_at",
    "SERMONS: id, title, speaker, series_id, video_url, thumbnail",
    "EVENTS: id, title, date, location, registration_url",
    "DONATIONS: id, user_id, amount, payment_method, status",
    "EVENT_REGISTRATIONS: id, event_id, user_id, status",
    "BOOKMARKS: id, user_id, sermon_id, created_at"
])

# Slide 12: Technology Stack
add_two_column_slide(prs, "Technology Stack",
    "Frontend (Complete):", [
        "HTML5 / CSS3 / JavaScript",
        "Tailwind CSS",
        "Google Fonts",
        "Heroicons / Font Awesome",
        "YouTube Embed",
        "Vercel Hosting"
    ],
    "Backend (To Develop):", [
        "Node.js / Express",
        "PostgreSQL Database",
        "Redis Caching",
        "JWT / OAuth 2.0",
        "Stripe / PayPal",
        "AWS S3 Storage"
    ])

# Slide 13: Design Principles
add_content_slide(prs, "Design Principles", [
    "MOBILE-FIRST: 78% of users access via mobile",
    "ACCESSIBILITY: WCAG 2.1 AA compliant",
    "BRAND CONSISTENCY: Navy, Gold, White palette",
    "PERFORMANCE: < 3 second load time",
    "SECURITY: HTTPS everywhere, input sanitization"
])

# Slide 14: Color Palette
add_content_slide(prs, "Brand Design System", [
    "PRIMARY: #0F172A (Navy) - Headers, backgrounds",
    "ACCENT: #3B82F6 (Blue) - CTAs, links",
    "GOLD: #D97706 - Highlights, accents",
    "LIGHT: #1E293B - Cards, secondary backgrounds",
    "Headings: Playfair Display (Serif)",
    "Body: Inter (Sans-serif)"
])

# Slide 15: Component Library
add_content_slide(prs, "Component Library", [
    "Navigation: Navbar, Mobile Hamburger Menu, Footer",
    "Cards: Sermon cards, Event cards, Staff cards",
    "Buttons: Primary, Secondary, Outline, Ghost",
    "Forms: Input fields, Validation, Submit buttons",
    "Modals: Login, Video Player, Confirmation",
    "50+ reusable components created"
])

# Slide 16: Responsive Design
add_content_slide(prs, "Responsive Breakpoints", [
    "MOBILE: < 640px - Single column, hamburger menu",
    "TABLET: 640-1024px - 2 columns, collapsed nav",
    "DESKTOP: > 1024px - Full navigation, 3-4 columns",
    "Touch targets: Minimum 44x44px for accessibility",
    "All pages tested across devices"
])

# Slide 17: Current Pages
add_content_slide(prs, "Pages Implemented (14 Total)", [
    "Home, About, Sermons, Events",
    "Contact, Give, Ministries",
    "Blog, Blog Post, Gallery",
    "Devotionals, Testimonies",
    "Privacy Policy, Terms of Service",
    "All with mobile hamburger menu"
])

# Slide 18: Version Roadmap
add_timeline_slide(prs, "Version Roadmap",
    {
        "v1.0 (Current)": ["MVP Launch", "14 Pages", "Basic Features"],
        "v1.1 (Q2)": ["User Auth", "Dashboard", "Bookmarks"],
        "v2.0 (Q3)": ["Payments", "Mobile App", "Analytics"],
        "v3.0 (Q4)": ["AI Features", "Multi-language", "Scale"]
    })

# Slide 19: 2026 Strategic Initiatives
add_content_slide(prs, "2026 Strategic Initiatives", [
    "Q1: Foundation - Complete frontend, begin backend",
    "Q2: Engagement - User auth, dashboard, events",
    "Q3: Growth - Payments, mobile app, analytics",
    "Q4: Scale - AI, multi-language, parish network"
])

# Slide 20: Value to ICT Department
add_content_slide(prs, "Adding Qualitative Value", [
    "REUSABLE COMPONENT LIBRARY - 40% faster future dev",
    "MODERN TECH STACK - Positions RCCG as tech-forward",
    "COMPREHENSIVE DOCS - Faster onboarding",
    "SCALABLE ARCHITECTURE - Future-proof growth",
    "KNOWLEDGE TRANSFER - Training & best practices"
])

# Slide 21: Team Structure
add_content_slide(prs, "Team Structure", [
    "ICT BOARD - Strategic decisions, budget",
    "PROJECT MANAGER - Sprint planning, timeline",
    "FRONTEND LEAD (You) - UI/UX, design, frontend",
    "BACKEND TEAM - API, database, auth",
    "DEVOPS - CI/CD, infrastructure, security"
])

# Slide 22: Your Role
add_content_slide(prs, "Role: Head of System Analysis & Frontend", [
    "SYSTEM ANALYSIS: Requirements, research, architecture",
    "FRONTEND DEVELOPMENT: UI/UX, components, performance",
    "PROJECT COORDINATION: Backend handoff, API specs, QA",
    "TEAM LEADERSHIP: Mentorship, code review, docs",
    "STAKEHOLDER MANAGEMENT: Board presentations"
])

# Slide 23: Budget Overview
add_content_slide(prs, "Budget Overview", [
    "Phase 1 (Complete): Frontend development, design, PM",
    "Phase 2 (Proposed): Backend, database, payments, testing",
    "Ongoing: Hosting, domain, maintenance",
    "ROI: Increased giving, engagement, efficiency"
])

# Slide 24: Risk Management
add_content_slide(prs, "Risk Management", [
    "SCOPE CREEP - Strict change management",
    "RESOURCE GAPS - Cross-training, contractors",
    "TECHNICAL DEBT - Regular refactoring",
    "INTEGRATION ISSUES - Early API definition",
    "USER ADOPTION - User testing, training"
])

# Slide 25: Achievements
add_content_slide(prs, "Achievements to Date", [
    "✓ Requirements analysis with 500+ survey responses",
    "✓ Complete wireframes for 14 pages",
    "✓ System architecture & database design",
    "✓ Frontend development complete",
    "✓ Vercel deployment configured",
    "✓ Documentation & onboarding guides"
])

# Slide 26: Next Steps
add_content_slide(prs, "Immediate Next Steps", [
    "BOARD APPROVAL: Backend development budget",
    "INFRASTRUCTURE: Cloud hosting setup",
    "TEAM EXPANSION: Onboard backend developers",
    "API DEVELOPMENT: Begin authentication system",
    "TIMELINE: Target Q2 2026 beta launch"
])

# Slide 27: KPIs
add_content_slide(prs, "Success Metrics & KPIs", [
    "TECHNICAL: < 3s load, 99.9% uptime, > 90 Lighthouse",
    "BUSINESS: 50K users, $500K donations, 85% retention",
    "OPERATIONAL: < 1hr content updates, weekly deploys",
    "USER: NPS > 50, 4.5+ rating, 90% satisfaction"
])

# Slide 28: Long-term Vision
add_content_slide(prs, "Long-term Vision 2027+", [
    "2027: Global parish network, AI personalization",
    "2028: Voice assistants, predictive analytics, VR",
    "2029: Industry leader, white-label, partner ecosystem",
    "ULTIMATE: Model for digital ministry globally"
])

# Slide 29: Q&A
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)
title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
title_bar.fill.solid()
title_bar.fill.fore_color.rgb = NAVY
title_bar.line.fill.background()

q_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(2))
tf = q_box.text_frame
p = tf.paragraphs[0]
p.text = "QUESTIONS & DISCUSSION"
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Slide 30: Thank You
add_closing_slide(prs, "THANK YOU",
    '"Building the Kingdom through Innovative Technology"',
    "Your Name\nHead of System Analysis & Frontend Design\nEmail: your.email@rccg.org\nGitHub: github.com/jdcrux1/RCCGMobileWorld")

# Save
output_file = "RCCG_MobileWorld_ICT_Board_Presentation.pptx"
prs.save(output_file)
print(f"Presentation saved as: {output_file}")
print("Done!")
